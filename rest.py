# import os
# import torch
# from dotenv import load_dotenv
# from langchain.prompts import PromptTemplate
# from langchain.memory import ConversationBufferMemory
# from langchain.chains import LLMChain
# from langchain.chains import SequentialChain
# from langchain import HuggingFaceHub
# import pptx
# from pptx.util import Pt
# import base64
# import streamlit as st
# from langchain_ollama import ChatOllama
# from langchain_core.prompts import ChatPromptTemplate

# load_dotenv()

# # Load Hugging Face API token
# os.environ['HUGGINGFACEHUB_API_TOKEN'] = os.getenv("HUGGINGFACEHUB_API_TOKEN")


# def generate_slide_titles(topic):
#     llm = ChatOllama(model="llama3.2", temperature=0)
#     prompt = ChatPromptTemplate.from_messages(
#         [
#             (
#                 "system",
#                 "Generate the one liner titles for the 5 page presentation"
#             ),
#             ("human", "{input}")
#         ]
#     )
    
#     chain = prompt | llm
#     res = chain.invoke(
#         {
#             "input": topic
#         }
#     )
    
#     return res



# def generate_content(name):
#     # Initialize the Ollama model with the correct parameters
#     llm = ChatOllama(model="llama3.2", temperature=0)
#     prompt = ChatPromptTemplate.from_messages(
#         [
#             (
#                 "system",
#                 "Write a structured 5-slide presentation outline, with content for each slide. Each slide should be detailed and focused on different aspects of the topic."
#             ),
#             ("human", "{input}")
#         ]
#     )
    
#     chain = prompt | llm
#     res = chain.invoke(
#         {
#             "input": name
#         }
#     )
    
#     # Process response to split into separate slide content
#     if isinstance(res.content, str):
#         content = res.content.split("\nSlide ")
#         content = [f"Slide {slide.strip()}" for slide in content if slide.strip()]
#     else:
#         content = []

#     return content

# def create_presentation(topic, slide_titles, slide_contents):
#     # Create a presentation object
#     prs = pptx.Presentation()

#     slide_layout = prs.slide_layout[1]
#     title_slide = prs.slides.add_slide(prs.slide_layouts[0])
#     title_slide.shapes.title.text = topic  # Use the topic as title
#     # title_slide.shapes.placeholders[1].text = "Presentation Overview"  # Set a generic subtitle

#     # Set the layout for the content slides
#     # slide_layout = prs.slide_layouts[1]

#     # Loop through the content list and create a slide for each
#     for slide_title , slide_content in zip(slide_titles, slide_contents):
#         slide = prs.slides.add_slide(slide_layout)  # Add a new slide
#         slide.shapes.title.text = slide_title
#         # Split slide title and content by the first colon
#         # slide_title, slide_text = slide_content.split(":", 1) if ":" in slide_content else (filename, slide_content)
        
#         # Set slide title and content
#         # slide.shapes.title.text = slide_title.strip()
#         content_placeholder = slide.shapes.placeholders[1]
#         content_placeholder.text = slide_content

#         # Format title text size
#         slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

#         # Format content text size
#         # for paragraph in content_placeholder.text_frame.paragraphs:
#         #     paragraph.font.size = Pt(18)  # Set the font size for each paragraph to 18 points
        
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text_frame = shape.text_frame
#                 for paragraph in text_frame.paragraphs:
#                     paragraph.font.size = Pt(18)
    
#     # Save the presentation
#     output_path = f"generated_ppt/{topic}_presentation.pptx"
#     os.makedirs(os.path.dirname(output_path), exist_ok=True)
#     prs.save(output_path)

# def main():
#     st.title("Generate PPT using LLM")
#     topic = st.text_input("Enter the topic")
#     generate_button = st.button("Generate PPT")

#     if generate_button and topic:
#         st.info("Generating PPT presentation.... Please wait.")
#         slide_titles = generate_slide_titles(topic)
        
#         filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
#         print("slide titles" , filtered_slide_titles)
        
        
#         # Generate content for each slide
#         slide_contents = [generate_content(topic) for titles in filtered_slide_titles]
#         print("slide content", slide_contents)
        
        
#         # Create the presentation
#         create_presentation(topic, filtered_slide_titles, slide_contents)
        
#         st.success("Presentation generated successfully ✅")
#         st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

# def get_ppt_download_link(topic_name):
#     ppt_path = f"generated_ppt/{topic_name}_presentation.pptx"
    
#     with open(ppt_path, "rb") as file:
#         ppt_contents = file.read()
    
#     b64_ppt = base64.b64encode(ppt_contents).decode()
    
#     return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download={ppt_path}>Download the PowerPoint file here</a>'

# if __name__ == '__main__':
#     main()

import os
import torch
from dotenv import load_dotenv
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
from langchain.chains import LLMChain
from langchain.chains import SequentialChain
import pptx
from pptx.util import Pt
import base64
import streamlit as st
from langchain_ollama import ChatOllama
from langchain_core.prompts import ChatPromptTemplate

load_dotenv()

# Load Hugging Face API token
os.environ['HUGGINGFACEHUB_API_TOKEN'] = os.getenv("HUGGINGFACEHUB_API_TOKEN")


def generate_slide_titles(topic):
    llm = ChatOllama(model="llama3.2", temperature=0)
    prompt = ChatPromptTemplate.from_messages(
        [
            (
                "system",
                "Generate the one-liner titles for a 5-page presentation."
            ),
            ("human", "{input}")
        ]
    )
    
    chain = prompt | llm
    res = chain.invoke({"input": topic})
    
    # Ensure res.content is a string and split into individual titles
    if isinstance(res.content, str):
        titles = res.content.split("\n")
    elif isinstance(res.content, tuple):
        titles = list(res.content)
    else:
        titles = []

    # Filter out any non-string values
    return [title for title in titles if isinstance(title, str) and title]


def generate_content(name):
    # Initialize the Ollama model with the correct parameters
    llm = ChatOllama(model="llama3.2", temperature=0)
    prompt = ChatPromptTemplate.from_messages(
        [
            (
                "system",
                "Write a structured 5-slide presentation outline, with content for each slide. Each slide should be detailed and focused on different aspects of the topic."
            ),
            ("human", "{input}")
        ]
    )
    
    chain = prompt | llm
    res = chain.invoke({"input": name})
    
    # Process response to split into separate slide content
    if isinstance(res.content, str):
        content = res.content.split("\nSlide ")
    elif isinstance(res.content, tuple):
        content = list(res.content)
    else:
        content = []

    # Filter out any non-string values
    return [f"Slide {slide}" for slide in content if isinstance(slide, str) and slide]


def create_presentation(topic, slide_titles, slide_contents):
    # Create a presentation object
    prs = pptx.Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    slide_layout = prs.slide_layouts[2]

    # Loop through the content list and create a slide for each
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title

        content_placeholder = slide.shapes.placeholders[1]
        content_placeholder.text = slide_content if isinstance(slide_content, str) else ""

        # Set font sizes
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)
        for paragraph in content_placeholder.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
    
    # Save the presentation
    output_path = f"generated_ppt/{topic}_presentation.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

def main():
    st.title("Generate PPT using LLM")
    topic = st.text_input("Enter the topic")
    generate_button = st.button("Generate PPT")

    if generate_button and topic:
        st.info("Generating PPT presentation... Please wait.")
        
        # Generate slide titles and content
        slide_titles = generate_slide_titles(topic)
        slide_contents = generate_content(topic)
        
        # Check that we have the right number of titles and contents
        if len(slide_titles) == 5 and len(slide_contents) == 5:
            create_presentation(topic, slide_titles, slide_contents)
            st.success("Presentation generated successfully ✅")
            st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)
        else:
            st.error("Could not generate the correct number of slides. Please try again.")

def get_ppt_download_link(topic_name):
    ppt_path = f"generated_ppt/{topic_name}_presentation.pptx"
    
    with open(ppt_path, "rb") as file:
        ppt_contents = file.read()
    
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{topic_name}_presentation.pptx">Download the PowerPoint file here</a>'

if __name__ == '__main__':
    main()

