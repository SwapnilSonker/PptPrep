import os
import torch
from diffusers import StableDiffusionPipeline
from dotenv import load_dotenv
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
from langchain.chains import LLMChain
from langchain.chains import SequentialChain
from langchain import HuggingFaceHub
import pptx
from pptx.util import Pt
import base64
import streamlit as st


load_dotenv()

os.environ['HUGGINGFACEHUB_API_TOKEN'] = os.getenv("HUGGINGFACEHUB_API_TOKEN")

def generate_content(name):
    
    person_memory = ConversationBufferMemory(input_key="name" , memory_key="chat_history")
    slide_memory = ConversationBufferMemory(input_key="n" , memory_key="chat_history")
    first_input_prompt = PromptTemplate(
    input_variables=['n', 'name'],
    template="Generate a {n} slide detailed content on the topic {name}"
    )

    llm = HuggingFaceHub(repo_id="google/flan-t5-large", model_kwargs={"temperature": 0.9, "max_length": 500})

    chain= LLMChain(llm=llm, prompt=first_input_prompt, verbose=True, output_key='person' , memory=person_memory)


    # name = "Artificial Intelligence"
    n=5
    res = chain.run({"name" : name, "n" :n })
    
    return res


def create_presentation(content, filename):
    # Create a presentation object
    prs = pptx.Presentation()

    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Presentation Title"  # Set a generic title
    title_slide.shapes.placeholders[1].text = "Subtitle"  # Set subtitle (optional)

    # Set the layout for the content slides
    slide_layout = prs.slide_layouts[1]

    # Loop through the content list and create a slide for each
    for slide_content in content:
        slide = prs.slides.add_slide(slide_layout)  # Add new slide
        slide.shapes.title.text = filename  # Placeholder for title
        content_placeholder = slide.shapes.placeholders[1]  # Placeholder for content
        
        # Set slide content
        content_placeholder.text = slide_content

        # Format title text size
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

        # Format content text size
        for paragraph in content_placeholder.text_frame.paragraphs:
            paragraph.font.size = Pt(18)  # Set the font size for each paragraph to 18 points
    return prs.save(f"generated_ppt/{filename}_presentation.pptx")     

def main():
    st.title("Generate PPT using LLM")
    topic = st.text_input("Enter the topic")
    generate_button = st.button("Generate PPT")

    if generate_button and topic:
        st.info("Generating PPT presentation.... Please wait.")
        slide_content = generate_content(topic).split(".")
        create_presentation(slide_content, topic)
        print("Presentation generated successfully...")
        
        st.success("Presentation generated successfully✅")
        st.markdown(get_ppt_download_link(topic) , unsafe_allow_html=True)

def get_ppt_download_link(topic_name):
    
    ppt_path = f"generated_ppt/{topic_name}_presentation.pptx"      
    
    with open(ppt_path,"rb") as file:
        ppt_contents = file.read()
        
    b64_ppt = base64.b64encode(ppt_contents).decode()
    
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download={ppt_path}>Download the powerpoint file here</a>'   

if __name__ == '__main__':
    main()   