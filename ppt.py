  

# import os
# import torch
# from diffusers import StableDiffusionPipeline
# from dotenv import load_dotenv
# from langchain.prompts import PromptTemplate
# from langchain.memory import ConversationBufferMemory
# from langchain.chains import LLMChain
# from langchain.chains import SequentialChain
# from langchain import HuggingFaceHub
# import pptx
# from pptx.util import Pt
# import base64
# import streamlit as st
# from langchain_ollama import ChatOllama
# from langchain_core.prompts import ChatPromptTemplate

# load_dotenv()

# os.environ['HUGGINGFACEHUB_API_TOKEN'] = os.getenv("HUGGINGFACEHUB_API_TOKEN")

# def generate_content(name):
#     # Update only the Ollama model section and ensure output is formatted as a list of slides
#     llm = ChatOllama(model="llama3.2", temperature=0)  # Correct model name and temperature format
#     prompt = ChatPromptTemplate.from_messages(
#         [
#             (
#                 "system",
#                 "Write me a detailed 5-page presentation, it should be especially crafted and detailed."
#             ),
#             ("human", "{input}")
#         ]
#     )
    
#     chain = prompt | llm
#     res = chain.invoke(
#         {
#             "input": name
#         }
#     )
    
#     # Ensure output is a list of strings (one per slide)
#     return res.content

# def create_presentation(content, filename):
#     # Create a presentation object
#     prs = pptx.Presentation()

#     # Add a title slide
#     title_slide = prs.slides.add_slide(prs.slide_layouts[0])
#     title_slide.shapes.title.text = "Presentation Title"  # Set a generic title
#     title_slide.shapes.placeholders[1].text = "Subtitle"  # Set subtitle (optional)

#     # Set the layout for the content slides
#     slide_layout = prs.slide_layouts[1]

#     # Loop through the content list and create a slide for each
#     for slide_content in content:
#         slide = prs.slides.add_slide(slide_layout)  # Add new slide
#         slide.shapes.title.text = filename  # Placeholder for title
#         content_placeholder = slide.shapes.placeholders[1]  # Placeholder for content
        
#         # Set slide content
#         content_placeholder.text = slide_content

#         # Format title text size
#         slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

#         # Format content text size
#         for paragraph in content_placeholder.text_frame.paragraphs:
#             paragraph.font.size = Pt(18)  # Set the font size for each paragraph to 18 points
#     return prs.save(f"generated_ppt/{filename}_presentation.pptx")     

# def main():
#     st.title("Generate PPT using LLM")
#     topic = st.text_input("Enter the topic")
#     generate_button = st.button("Generate PPT")

#     if generate_button and topic:
#         st.info("Generating PPT presentation.... Please wait.")
#         slide_content = generate_content(topic)
#         create_presentation(slide_content, topic)
#         print("Presentation generated successfully...")
        
#         st.success("Presentation generated successfully✅")
#         st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

# def get_ppt_download_link(topic_name):
#     ppt_path = f"generated_ppt/{topic_name}_presentation.pptx"      
    
#     with open(ppt_path, "rb") as file:
#         ppt_contents = file.read()
        
#     b64_ppt = base64.b64encode(ppt_contents).decode()
    
#     return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download={ppt_path}>Download the powerpoint file here</a>'   

# if __name__ == '__main__':
#     main()

import os
import torch
# from diffusers import StableDiffusionPipeline
from dotenv import load_dotenv
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
from langchain.chains import LLMChain
from langchain.chains import SequentialChain
from langchain import HuggingFaceHub
import pptx
from pptx.util import Pt
import base64
import streamlit as st
from langchain_ollama import ChatOllama
from langchain_core.prompts import ChatPromptTemplate

load_dotenv()

os.environ['HUGGINGFACEHUB_API_TOKEN'] = os.getenv("HUGGINGFACEHUB_API_TOKEN")

def generate_content(name):
    # Generate structured presentation content using the LLM
    llm = ChatOllama(model="llama3.2", temperature=0.7)  # Correct model name and temperature format
    prompt = ChatPromptTemplate.from_messages(
        [
            (
                "system",
                "Write a structured 5-slide presentation outline, with content for each slide. Each slide must generate elaborated bullet points with no images included on the particular point."
            ),
            ("human", "{input}")
        ]
    )
    
    chain = prompt | llm
    res = chain.invoke({"input": name})
    
    # Split response content into slide sections, assuming new slide content starts with "Slide X:"
    content = res.content.split("\nSlide ") if isinstance(res.content, str) else []
    content = [slide.strip() for slide in content if slide]  # Clean up whitespace
    return content

def create_presentation(content, filename):
    # Create a presentation object
    prs = pptx.Presentation()

    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = filename  # Use the topic as title
    title_slide.shapes.placeholders[1].text = "Presentation Overview"  # Set a generic subtitle

    # Set layout for content slides
    slide_layout = prs.slide_layouts[1]  # Title and Content layout

    # Loop through each slide's content and create a new slide
    for i, slide_content in enumerate(content):
        # Split each slide content into title and body if separated by a colon
        slide_title, slide_text = (slide_content.split(":", 1) if ":" in slide_content 
                                   else (f"Slide {i + 1}", slide_content))
        
        # Add new slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title and content
        slide.shapes.title.text = slide_title.strip()
        content_placeholder = slide.shapes.placeholders[1]
        content_placeholder.text = slide_text.strip()

        # Set font size for title and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)
        for paragraph in content_placeholder.text_frame.paragraphs:
            paragraph.font.size = Pt(18)  # Set content font size

    # Save the presentation
    output_path = f"generated_ppt/{filename}_presentation.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved as {output_path}")
def main():
    st.title("Generate PPT using LLM")
    topic = st.text_input("Enter the topic")
    generate_button = st.button("Generate PPT")

    if generate_button and topic:
        st.info("Generating PPT presentation.... Please wait.")
        slide_content = generate_content(topic)
        create_presentation(slide_content, topic)
        print("Presentation generated successfully...")
        
        st.success("Presentation generated successfully✅")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

def get_ppt_download_link(topic_name):
    ppt_path = f"generated_ppt/{topic_name}_presentation.pptx"      
    
    with open(ppt_path, "rb") as file:
        ppt_contents = file.read()
        
    b64_ppt = base64.b64encode(ppt_contents).decode()
    
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download={ppt_path}>Download the powerpoint file here</a>'   

if __name__ == '__main__':
    main()
